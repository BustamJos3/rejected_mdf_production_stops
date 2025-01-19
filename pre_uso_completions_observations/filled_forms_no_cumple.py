#!/usr/bin/env python
# coding: utf-8

# In[15]:


import pandas as pd #modules
import numpy as np
import matplotlib
import matplotlib.pyplot as plt
import matplotlib.colors as mcolors
import os
import pathlib
from pathlib import Path
from datetime import datetime, timedelta
import itertools
import string


# # Load shift hours

# In[16]:


def find_files_by_format(root_folder, file_format):
    """
    Searches for all files with a specific format in a folder and its subfolders.

    Parameters:
    root_folder (str): The root folder to start the search.
    file_format (str): The file extension to search for (e.g., '.pptx', '.txt').

    Returns:
    list: A list of absolute paths to all matching files.
    """
    matching_files = []

    # Traverse directories recursively
    for root, dirs, files in os.walk(root_folder):
        for file in files:
            # Check if the file has the specified format
            if file.lower().endswith(file_format.lower()):  # Case-insensitive check
                # Append the full path of the file
                matching_files.append(os.path.join(root, file))

    return matching_files


# In[17]:


def operate_mercyful_times(hour,operand): #function to apply modifications on shift hours
    from datetime import datetime, timedelta
    delta_time_mercy=timedelta(minutes=0)
    mercy_flag=False #bool to activate mercy mode
    if mercy_flag:
        delta_time_mercy=timedelta(minutes=10) #time delta to substract and add to initial and final hours
    if operand=="+" and mercy_flag: #to add to final hour
        result_hour=eval("datetime.strptime(hour, '%H:%M:%S'){}delta_time_mercy".format(operand))
        result_hour=result_hour.time()
        return result_hour
    elif operand=="-" and mercy_flag: #to substract to initial hour
        result_hour=eval("datetime.strptime(hour, '%H:%M:%S'){}delta_time_mercy".format(operand))
        result_hour=result_hour.time()
        return result_hour
    else:
        return eval("datetime.strptime(hour, '%H:%M:%S').time()")


# In[18]:


str___splitter=lambda x: x.split("_")


# # Function to find cloud (OneDrive) folder

# In[19]:


import os

def find_reports_in_onedrive():
    """
    Scans the subfolders under the current user's OneDrive folder (including variations like 'OneDrive - Company Name')
    and returns the paths of all folders with the prefix 'reports'.

    Returns:
        list: A list of full paths to folders starting with 'reports', or an empty list if none are found.
    """
    # Get the base path to the user's home directory
    user_home = os.path.expanduser("~")

    # Find the OneDrive folder (handles variations like "OneDrive - Company Name")
    onedrive_folder = None
    for folder in os.listdir(user_home):
        if folder.startswith("OneDrive -"):
            onedrive_folder = os.path.join(user_home, folder)
            break

    if not onedrive_folder:
        raise FileNotFoundError("OneDrive folder not found for the current user.")

    # Search for folders with the prefix 'reports' in the OneDrive directory
    report_folders = []
    for root, dirs, files in os.walk(onedrive_folder):
        for dir_name in dirs:
            if dir_name.lower().startswith("reports"):
                report_folders.append(os.path.join(root, dir_name))

    return report_folders


# In[20]:


reports_paths=find_reports_in_onedrive()
reports_paths


# In[21]:


str_folder_searcher="reports_pre-usos_diligenciamiento_no_cumple"
for report_path in reports_paths:
    if str_folder_searcher in report_path:
        path=Path(reports_paths[reports_paths.index(report_path)])
path=Path.joinpath(path,r"source_and_return_data")
path


# In[22]:


os.chdir(path) #change working directory to respective reports folder
os.getcwd()


# In[23]:


shifts_hours=pd.read_csv(r".\shifts_hours.txt",delimiter=",")
shifts_hours_cols=list(shifts_hours.columns)
shifts_hours.loc[:,shifts_hours_cols[0]]=shifts_hours.loc[:,shifts_hours_cols[0]].apply(lambda x: operate_mercyful_times(hour=x,operand="-"))
shifts_hours.loc[:,shifts_hours_cols[1]]=shifts_hours.loc[:,shifts_hours_cols[1]].apply(lambda x: operate_mercyful_times(hour=x,operand="+"))
shifts_hours.loc[:,shifts_hours_cols[-1]]=shifts_hours.loc[:,shifts_hours_cols[-1]].apply(str___splitter)
shifts_hours


# # Generation of querie to filter dates to look up completion dates

# In[24]:


def get_monday_of_week():
    """
    This function, `get_monday_of_week`, returns the date  
    of the Monday in the current week based on today's date.  
    
    It first retrieves today's date using `datetime.today()`  
    and then determines the current weekday using  
    `today.weekday()`, where Monday is represented by 0.  
    
    Next, it calculates the number of days to subtract to  
    reach Monday by taking the value of `today.weekday()`  
    and subtracting that number of days using `timedelta`.  
    
    Finally, it formats the result as 'YYYY-MM-DD' and  
    returns it as the output.  
    """
    # Obtener la fecha de hoy
    today = datetime.today()
    # Calcular la diferencia de días hasta el lunes (Monday = 0)
    delta = today.weekday()  # weekday() devuelve 0 para lunes
    # Restar los días necesarios para llegar al lunes
    monday = today - timedelta(days=delta)
    # Retornar la fecha en formato YYYY-MM-DD
    return monday.date()


# In[25]:


def date_limits(moment_of_week=True,day_separation=1):
    """change dates limits to filter completions 
    base on bool (default True), whether is weekend 
    or in between the week and if the moment in between 
    the week is 1 day apart or more"""
    if not moment_of_week: #day of querie is not in between the week
        day_separation=3 #querie to consult sunday (1), saturday (2) and friday (3)
        get_current_recent_date=get_monday_of_week() #get the date of monday from current week
    else:
        get_current_recent_date=datetime.now()
        print(get_current_recent_date)
    day_to_stop_search=get_current_recent_date-timedelta(days=1) #check from yesterday
    print(day_to_stop_search)
    day_before_yesterday=day_to_stop_search - timedelta(days=day_separation) #get yesterday date to check fraction of T1 and T2 y T3
    print(day_before_yesterday)
    return day_to_stop_search,day_before_yesterday


# ## Construct separation of days from last execution date to day to search in pre-usos
# * store last execution date to compute separation of days

# In[14]:


today=datetime.now() #get date_time from today
today_date=today.date() #get date
# Specify the file name
file_name = r".\date_time_last_exec.txt"
if today_date==get_monday_of_week(): #if it is monday
    not_weekend=False #moment of querie is weekend
    separation_of_days=None
else:
    not_weekend=True #moment of querie is not weekend
    # Read the content of the file and parse it as a datetime object
    with open(file_name, 'r') as file:
        last_execution = datetime.strptime(file.read().strip(), '%Y-%m-%d %H:%M:%S') #load last execution
        last_execution_date=last_execution.date()
        print(last_execution_date)
    separation_of_days=today_date-last_execution_date
    separation_of_days=separation_of_days.days if separation_of_days.days>0 else 365+separation_of_days.days #days separation to make queries in between week of more than 1 day
    print("separation",separation_of_days)
"""export last execution date into txt"""
# Write the date_time variable into the file
list_current_reports=find_files_by_format(Path.joinpath(path,"pptx_reports"),".pptx")
yesterday=today_date-timedelta(days=1)
most_recent_report_title_yesterday=False #flag to check yesterday report existance
for report_title in list_current_reports:
    if yesterday.strftime('%Y-%m-%d') in report_title:
        most_recent_report_title_yesterday=True
if most_recent_report_title_not_yesterday: #if a report with yesterday date is not already generated, update last execution
    with open(file_name, 'w') as file:
        print("no matches in reports, proceed to update last execution")
        file.write(today.strftime('%Y-%m-%d %H:%M:%S'))
day_to_stop_search_querie,day_before_yesterday_querie=date_limits(moment_of_week=not_weekend,day_separation=separation_of_days)
day_before_yesterday_querie=datetime.combine(day_before_yesterday_querie,shifts_hours.at[0,"initial_hour"])
day_to_stop_search_querie=datetime.combine(day_to_stop_search_querie,shifts_hours.at[2,"finalitation_hour"])
day_to_stop_search_querie,day_before_yesterday_querie


# # Load pre-usos files

# In[13]:


directory = Path(r".\pre_usos_actual") #get current work directory
matching_files = list(directory.glob("*Pre-Uso*.xlsx"))  # Busca archivos que tengan extensión .xlsx
print("Archivos encontrados:", matching_files)


# In[14]:


dict_data_pointer={} #dict to store files as dfs
for i in matching_files: # Read the Excel file 
    file_path = str(i)  # Update this with the path
    try:
        df = pd.read_excel(file_path,sheet_name="Form1") #specific sheet call for english origin forms
    except:
        df = pd.read_excel(file_path)
    print(file_path)
    df_name=file_path.split(str(directory)+"\\")[-1].split(".")[0] #split str with "." char and take the file name
    dict_data_pointer[f"{df_name}"]=df #store filtered df in dictionary data pointer
print(list(dict_data_pointer.keys())) #see keys on dictionary to check callability


# # Get forms done respect the shifts

# # Make df to create excel with table of completions and specific rules of laboral days

# In[15]:


day_difference=day_to_stop_search_querie-day_before_yesterday_querie
day_difference=day_difference.days
day_difference


# ## change creation of row_report with the reading from loading the actual workbook with the dates

# In[16]:


df_pre_uso_completion_report=pd.read_excel(Path.joinpath(path,r"pre_uso_produccion_report_copia\Preusos Producción.xlsx"))
df_pre_uso_completion_report_cols=list(df_pre_uso_completion_report.columns)
array_data=[["","",""]*day_difference]*len(df_pre_uso_completion_report_cols[2:])
array_data=list(map(list, zip(*array_data))) #transpose list to match df input data format
print(array_data)
shifts=3
dates_of_interest= [ [date]*shifts for date in [day_to_stop_search_querie.date() - timedelta(days=x) for x in range(day_difference)] ]
dates_of_interest=list(itertools.chain.from_iterable(dates_of_interest))
if len(dates_of_interest)>1:
    dates_of_interest=sorted(dates_of_interest)
multi_index=[ dates_of_interest,["1","2","3"]*day_difference ]
print(multi_index)
df_check_completion=pd.DataFrame(data=array_data,index=multi_index,columns=df_pre_uso_completion_report_cols[2:])
df_check_completion.index.names =['Día',df_pre_uso_completion_report_cols[1]]
df_check_completion


# ## Function to check shifts

# In[17]:


def shift_checker(date_time=None,system_to_test=None):
    system=system_to_test.lower() #convert system name to lowercase
    print(date_time)
    time_of_interest=date_time.time() #take the hh:mm:ss of entry in completion date_time
    initial_shift_hours=shifts_hours.loc[:,[shifts_hours_cols[0]]].values
    final_shift_hours=shifts_hours.loc[:,[shifts_hours_cols[1]]].values
    available_systems=[sys[0] for sys in shifts_hours.loc[:,[shifts_hours_cols[-1]]].values]
    for available_system in available_systems: #check which is the shift of the system
        if system in available_system:
            print(available_system)
            matching_shift=available_systems.index(available_system)+1
            print(matching_shift)
            if matching_shift<4: #system is from shift 1,2 or 3
                print(system," in ",matching_shift)
                if time_of_interest>=initial_shift_hours[0][0] or time_of_interest<final_shift_hours[0][0]: #check if it is shift 1
                    shift_result="T1"
                if time_of_interest>=initial_shift_hours[1][0] and time_of_interest<final_shift_hours[1][0]: #check if it is shift 2
                    shift_result="T2"
                if time_of_interest>=initial_shift_hours[2][0] and time_of_interest<final_shift_hours[2][0]: #check if it is shift 3
                    shift_result="T3"
            else: #system if from shift 4
                if time_of_interest>=initial_shift_hours[3][0] and time_of_interest<final_shift_hours[3][0]: #check if it is shift 3
                    shift_result="T4"
            return shift_result
        else:
            print("sistema no reconocido en {}".format(available_system))


# ## Apply check of shifts on dfs pre-usos

# In[18]:


dict_result_shift_checker={}
for pre_uso in dict_data_pointer.keys(): #run throug all dfs of pre-usos
    df=dict_data_pointer[pre_uso] #select respective df
    df_columns=list(df.columns)
    df.loc[:,df_columns[1]] = pd.to_datetime(df.loc[:,df_columns[1]], format='%Y-%m-%d %H:%M:%S') # Convierte las columnas al formato datetime
    df_shift_dates=df.loc[( df[df_columns[1]]>=day_before_yesterday_querie )&( df[df_columns[1]]<=day_to_stop_search_querie )] #filter df to get completions of desired date
    df_shift_dates=df_shift_dates.sort_values(by=df_columns[1],ascending=False)
    df_shift_dates=df_shift_dates.fillna(value=" ")
    df_np_times=df_shift_dates.loc[:,df_columns[1]].values #get the dates of current df filtered
    """start shift check"""
    system_to_check_in=pre_uso.lower().split("pre-uso ")[-1]
    print(f"sistema a chequear: {system_to_check_in}")
    dict_result_shift_checker[system_to_check_in]=[]
    for date_np_time in df_np_times:
        date_timestamp=pd.Timestamp(date_np_time)
        print(f"fecha de diligenciamiento: {date_timestamp}")
        shift_checked=shift_checker(date_time=date_timestamp,system_to_test=system_to_check_in) #store shift
        print(f"turno asignado: {shift_checked}")
        dict_result_shift_checker[system_to_check_in].append((date_timestamp,shift_checked))
    dict_result_shift_checker[system_to_check_in]=sorted(dict_result_shift_checker[system_to_check_in])
dict_result_shift_checker


# ## asign values on df of completion

# In[19]:


df_check_completion_cols=df_check_completion.columns
for (key_shift_checker,values_shift_checker) in dict_result_shift_checker.items():
    print(key_shift_checker,values_shift_checker)
    for current_col_check_completion in df_check_completion_cols:
        if key_shift_checker in current_col_check_completion.lower():
            print(f"sistema {key_shift_checker} in {current_col_check_completion}")
            for date_time,shift in values_shift_checker: #begin asignation by dates and shifts
                current_date_of_shift_checker=date_time.date()
                idx_shift=shift[-1]
                print(current_date_of_shift_checker)
                if current_date_of_shift_checker not in dates_of_interest:
                    current_date_of_shift_checker=min(dates_of_interest) #if date not in dates_of_interest, take the youngest date
                print(idx_shift)
                df_check_completion.loc[(current_date_of_shift_checker, idx_shift), current_col_check_completion]="OK" #asignation
df_check_completion


# ## export completion df to excel

# In[20]:


path_to_export=Path.joinpath(path,rf'row_date_report\{day_to_stop_search_querie.date()}.xlsx')
print(path_to_export)
df_check_completion.to_excel(path_to_export)


# # Filter dfs of preusos on one cell to get "No Cumple"

# In[21]:


dict_data_no_cumple={}
for pre_uso_key in list(dict_data_pointer.keys()):
    print(pre_uso_key)
    """select df"""
    df=dict_data_pointer[pre_uso_key]
    """save columns names"""
    df_columns=list(df.columns)
    # Convierte las columnas al formato datetime
    df.loc[:,df_columns[1]] = pd.to_datetime(df.loc[:,df_columns[1]], format='%Y-%m-%d %H:%M:%S')
    """filter dates to check completion"""
    df_shift_dates=df.loc[( df[df_columns[1]]>=day_before_yesterday_querie )&( df[df_columns[1]]<=day_to_stop_search_querie )] #get completion dates
    df_shift_dates=df_shift_dates.sort_values(by=df_columns[1],ascending=False)
    df_shift_dates=df_shift_dates.fillna(value=" ")
    
    """create df to check observations on No Cumple entries"""
    df_to_concatenate=pd.DataFrame() #df to begin concatenation
    for col_name in df_columns[8:]:
        #print(col_name)
        print(df_shift_dates[col_name])
        #try: #try to get No completion entries
        df_col_name_querie=df_shift_dates.loc[ (df_shift_dates[col_name].astype("str").str.contains("No Cumple"))|(df_shift_dates[col_name].astype("str").str.contains("No cumple")) ] #--->reformat to make cleaner solution
        #except: #continue cause there is no entries with No Cumple
            #continue
        if len(df_col_name_querie)>0: #"No Cumple" is found in current col
            print(col_name)
            df_with_concatenate=df_col_name_querie.loc[:,[col_name]]#df_col_name_querie.loc[:,[df_columns[1],df_columns[-1],col_name]]
            df_to_concatenate=pd.concat([df_to_concatenate, df_with_concatenate],axis=1, sort=False)
    indexes_dates_observs=list(df_to_concatenate.index)
    df_to_concatenate=pd.concat([df_shift_dates.loc[indexes_dates_observs,[df_columns[-1],df_columns[1] ]],df_to_concatenate],axis=1, sort=False)
    df_to_concatenate=df_to_concatenate.fillna(value=" ")#sort_values(by=df_columns[1])
    if len(df_to_concatenate)>0: #only take the queries that are not empty
        dict_data_no_cumple[pre_uso_key]=df_to_concatenate


# # Construction of data to put on presentation

# # Write in presentation
# ## TODO
# * If size of observations is adequate, put no completion entries, questions label and observations in the same slide
# * Table must fit on one slide without exceding the size of the slide

# In[22]:


from pptx import Presentation
from pptx.util import Inches, Pt


# In[23]:


def iter_cells(table):
    """
    to change font of text in table
    """
    for row in table.rows:
        for cell in row.cells:
            yield cell


# In[24]:


prs = Presentation()
"""make title slide"""
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Pre-Usos BAR"
subtitle.text = "{}, {}".format(day_before_yesterday_querie.date(),day_to_stop_search_querie.date()) #datetime.now().date()
"""make next slides to inform about pre-usos"""
"""make completion report slide"""
title_no_content_slide_layout=prs.slide_layouts[5]
slide = prs.slides.add_slide(title_no_content_slide_layout)
title = slide.shapes.title
title.text=f"Pre-Usos periodo {subtitle.text}"
title_para = slide.shapes.title.text_frame.paragraphs[0]
title_para.font.size = Pt(20) #set title size
"""make slides to report about No Cumple"""
x, y, cx, cy = Inches(0.1), Inches(1.2), Inches(6), Inches(5)
factorx=0.4
factory=0.25
offsets=[cx*(1-factorx),cy*(1+factory)] #to adjust textbox @right of table
simple_content_slide_layout=prs.slide_layouts[1]
max_font_to_fit_text=16 #value for max font size when fitting text
for pre_uso_name in dict_data_no_cumple.keys(): #take each df with no cumple entries
    slide=prs.slides.add_slide(title_no_content_slide_layout) #slide to add dates and no completion items
    title = slide.shapes.title
    title.text=pre_uso_name

    title_para = slide.shapes.title.text_frame.paragraphs[0]
    title_para.font.size = Pt(20) #set title size
    
    df=dict_data_no_cumple[pre_uso_name].T
    table_rows,table_cols=df.shape
    #print(table_rows,table_cols)
    """insert text box label to identify items with questions"""
    txBox=slide.shapes.add_textbox(x+cx, y, offsets[0], offsets[1])
    tf = txBox.text_frame
    tf.word_wrap = True
    #tf.autofit_text()
    p = tf.add_paragraph()
    text_container = "" #var to store questions of respective pre-uso
    #p.font.size = Pt(12)
    """make table of no completions"""
    shape = slide.shapes.add_table(table_rows-1, table_cols+1, x, y, cx, cy) #table geometry
    table=shape.table
    cell = table.cell(0, 0)
    cell.text="Item"
    df_col_index=list(df.index)[1:] #name of current index which were columns before
    print(df_col_index)
    df_index_col=list(df.columns) #name of current cols which were indexes before
    for idx_col in range(table_cols): #write on header cells of table
        print("indice de columna en tabla: ",idx_col)
        cell = table.cell(0, idx_col+1)
        cell.text="{}".format(df.at[ df_col_index[0],df_index_col[idx_col] ])
     #write on first column to put questions with no completion
    for idx_row in range(table_rows-2):
        print("indice de fila en tabla: ",idx_row)
        cell = table.cell(idx_row+1, 0)
        print(df_col_index[idx_row+1])
        cell.text="{}".format(idx_row+1) #df_col_index[idx_row+1]
        if idx_row<1:
            text_container+= f"->Pregunta {idx_row+1}: "+df_col_index[idx_row+1] #insert text to label_id_questions text box without next line insertion
        else:
            text_container+= f"\n->Pregunta {idx_row+1}: "+df_col_index[idx_row+1] #insert text to label_id_questions text box 
    p.text=text_container #put text on textbox
    tf.fit_text(max_size=max_font_to_fit_text) #adjust text to fit int
    #insert respective data
    for idx_row in range(table_rows-2):
        for idx_col in range(table_cols):
            print("celda value:",idx_row+1,idx_col)
            cell = table.cell(idx_row+1, idx_col+1)
            cell.text="{}".format(df.at[ df_col_index[idx_row+1],df_index_col[idx_col] ]) #insert df value
    """to change table text font"""
    for cell in iter_cells(table):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(12)
    """to add observations"""
    string_to_add_observations=""
    flag_make_observations_slide=False
    for idx_col in range(table_cols):
    #put date of observation and observation content
        observation=df.at[ list(df.index)[0],df_index_col[idx_col] ]
        min_lenght_observation=40 #min lenght of observation to be inserted on slides report
        if len(observation)<min_lenght_observation: #if observation lenght is minor to 28 chars, continue
            continue
        else:
            flag_make_observations_slide=True
            if idx_col<1:
                string_to_add_observations+=f"Observación {df.at[ df_col_index[0],df_index_col[idx_col] ]}: {observation}" #insert 1st observation without next line insertion
            else:
                string_to_add_observations+=f"\nObservación {df.at[ df_col_index[0],df_index_col[idx_col] ]}: {observation}" #insert 1st observation without next line insertion            
    if flag_make_observations_slide:
        slide=prs.slides.add_slide(simple_content_slide_layout) #slide to add observations
        title = slide.shapes.title
        title.text="{}-{}".format(pre_uso_name,"Observaciones")
    
        title_para = slide.shapes.title.text_frame.paragraphs[0]
        title_para.font.size = Pt(28) #set title size
        """insert content on text box place holder"""
        content_tf=slide.placeholders[1]
        format_text_props=content_tf.text_frame
        format_text_props.word_wrap=True
        content_tf.text=string_to_add_observations
        format_text_props.fit_text(max_size=max_font_to_fit_text) #adjust text to fit int

presentation_date=day_to_stop_search_querie.date()#.strftime('%d-%m-%Y')
print(presentation_date)
year_folder=rf"\year_{presentation_date.year}"
month_folder=rf"\month_{day_to_stop_search_querie.month}"
directory=r"pptx_reports"
directory_to_save = Path.joinpath(path,str(directory)+year_folder+month_folder) # get directory to save plot
directory_to_save.mkdir(exist_ok=True)
prs.save(str(directory_to_save)+rf"\Preusos Producción y MMTTO_{presentation_date}.pptx")
prs.save(Path.joinpath(path,str(directory)+r'\test\test.pptx'))


# In[1]:


get_ipython().system('jupyter nbconvert filled_forms_no_cumple.ipynb --to script')

